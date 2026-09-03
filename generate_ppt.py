import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(15, 23, 42)       # Slate 900
    COLOR_CARD = RGBColor(30, 41, 59)     # Slate 800
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_CYAN = RGBColor(14, 165, 233)    # Cyan 500
    COLOR_PURPLE = RGBColor(139, 92, 246)  # Purple 500
    COLOR_ROSE = RGBColor(244, 63, 94)    # Rose 500
    COLOR_AMBER = RGBColor(245, 158, 11)   # Amber 500
    COLOR_MUTED = RGBColor(148, 163, 184)  # Slate 400

    def add_blank_slide(bg_color=COLOR_BG):
        slide_layout = prs.slide_layouts[6] # blank
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        return slide

    def add_header(slide, tag_text, title_text, category_color=COLOR_CYAN):
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = category_color

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    def add_card(slide, left, top, width, height, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    s1 = add_blank_slide()
    # Brand Tag
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "HUMAN-MANIPULATION FRAUD FIREWALL FOR BANKS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    # Main Title
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "FinShield Quantum"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Subtitle
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.3), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "“Detecting the manipulation behind legitimate transactions.”"
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_MUTED

    # Highlights box
    add_card(s1, 1.0, 4.8, 11.33, 1.8, border_color=COLOR_CYAN)
    tb = s1.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10.9), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multimodal AI • Adaptive Firewall • Senior Safety (WCAG AAA) • Post-Quantum Cryptography"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Hackathon Pitch Deck | Enterprise Banking Security Infrastructure Prototype"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_MUTED
    p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 2: The Problem (APP Scams)
    # ==========================================
    s2 = add_blank_slide()
    add_header(s2, "Industry Threat Landscape", "The $1.5B+ Authorized Push Payment (APP) Scam Crisis")

    # Card 1: The Problem
    add_card(s2, 0.8, 1.7, 5.6, 5.0, border_color=COLOR_ROSE)
    tb = s2.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.2), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Why Traditional Fraud Detection Fails"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE

    bullets1 = [
        "APP Scams deceive victims into authorizing wire transfers to fraud accounts.",
        "Social Engineering Coercion: Impersonating bank officers, law enforcement, or urgent family distress.",
        "Technical Legitimacy: Victims use their real phone, correct password, and pass 2FA/MFA.",
        "Traditional Fraud Systems see: 'AUTHENTIC TRANSACTION' and allow funds to leave.",
        "Devastating Impact: Older adults and digitally inexperienced users lose lifetime savings."
    ]
    for b in bullets1:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(8)

    # Card 2: Key Stats
    add_card(s2, 6.8, 1.7, 5.7, 5.0, border_color=COLOR_CYAN)
    tb = s2.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.3), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Market Impact & Regulatory Mandate"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    stats = [
        ("$1.5 Billion+", "Global losses to Authorized Push Payment scams in 2025/2026."),
        ("72% Senior Target", "Older adults targeted by vishing coercion calls."),
        ("50% Bank Liability", "New UK PSR & US regulatory rules forcing banks to reimburse scam victims."),
        ("< 50ms Window", "Banks require real-time inline evaluation before payment finality.")
    ]
    for val, lbl in stats:
        p = tf.add_paragraph()
        p.text = val
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_AMBER
        p.space_before = Pt(10)
        
        p_sub = tf.add_paragraph()
        p_sub.text = lbl
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 3: The Critical Paradigm Shift
    # ==========================================
    s3 = add_blank_slide()
    add_header(s3, "FinShield Core Innovation", "The Critical Question Paradigm Shift")

    # Left: Old Question
    add_card(s3, 0.8, 1.8, 5.6, 4.8, border_color=COLOR_ROSE)
    tb = s3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "TRADITIONAL FRAUD SYSTEM"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED

    p = tf.add_paragraph()
    p.text = "“Is the transaction legitimate?”"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE
    p.space_before = Pt(10)

    checks1 = [
        "✓ Real User Device ID",
        "✓ Valid Password & PIN",
        "✓ Passed 2FA / OTP",
        "✓ Legitimate Session Token",
        "RESULT: TRANSACTION ALLOWED ❌ (Scam Succeeded)"
    ]
    for c in checks1:
        p = tf.add_paragraph()
        p.text = c
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(10)

    # Right: FinShield Question
    add_card(s3, 6.8, 1.8, 5.7, 4.8, border_color=COLOR_CYAN)
    tb = s3.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "FINSHIELD QUANTUM FIREWALL"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p = tf.add_paragraph()
    p.text = "“Is the CUSTOMER'S INTENT being manipulated?”"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_before = Pt(10)

    checks2 = [
        "⚠ Scans incoming SMS/calls for coercion NLP",
        "⚠ Evaluates user spending z-score anomaly",
        "⚠ Checks Payee network graph reputation",
        "⚠ Factors in user vulnerability context",
        "RESULT: TRANSACTION HELD ⛔ (₹85,000 Saved)"
    ]
    for c in checks2:
        p = tf.add_paragraph()
        p.text = c
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(10)

    # ==========================================
    # SLIDE 4: Multimodal AI Architecture
    # ==========================================
    s4 = add_blank_slide()
    add_header(s4, "Technical Architecture", "Multimodal Intelligence Engine Pipeline")

    sub_engines = [
        ("1. Transaction ML", "XGBoost Classifier evaluating amount ratio, location shift, and time of day.", COLOR_CYAN),
        ("2. Behavior Baseline", "Statistical Z-Score Anomaly Detector measuring distance from normal spending.", COLOR_PURPLE),
        ("3. Social NLP", "Language Scorer evaluating authority impersonation & urgency keywords in calls/SMS.", COLOR_ROSE),
        ("4. Fraud Graph", "NetworkX Graph computing Payee centrality, shared devices, and collusive scam rings.", COLOR_AMBER)
    ]

    for idx, (title, desc, col) in enumerate(sub_engines):
        left = 0.8 + (idx * 2.95)
        add_card(s4, left, 1.8, 2.8, 3.8, border_color=col)
        tb = s4.shapes.add_textbox(Inches(left + 0.15), Inches(2.0), Inches(2.5), Inches(3.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(10)

    # Fusion Bottom Banner
    add_card(s4, 0.8, 5.8, 11.7, 1.2, border_color=COLOR_CYAN)
    tb = s4.shapes.add_textbox(Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multimodal Risk Fusion Formula: Manipulation Score = 0.35(Social NLP) + 0.25(Txn ML) + 0.20(Behavior) + 0.20(Graph) [Adaptive Senior Boost]"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 5: 4-Tier Adaptive Firewall Policy
    # ==========================================
    s5 = add_blank_slide()
    add_header(s5, "Dynamic Policy Engine", "4-Tier Adaptive Firewall Matrix")

    tiers = [
        ("ALLOW", "0 - 29 Risk", "Green Pass", "Payment completes seamlessly with zero user friction.", COLOR_CYAN),
        ("VERIFY", "30 - 59 Risk", "Yellow 2FA", "Prompts biometrics or step-up OTP verification.", COLOR_AMBER),
        ("PAUSE", "60 - 84 Risk", "Orange Reflection", "Pauses transaction for 60 seconds with safety prompt.", COLOR_PURPLE),
        ("HOLD", "85 - 100 Risk", "Red Interception", "Blocks transfer, alerts trusted contact, requires bank call.", COLOR_ROSE)
    ]

    for idx, (t_name, t_range, t_tag, t_desc, col) in enumerate(tiers):
        top = 1.8 + (idx * 1.25)
        add_card(s5, 0.8, top, 11.7, 1.1, border_color=col)
        tb = s5.shapes.add_textbox(Inches(1.0), Inches(top + 0.15), Inches(11.3), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{t_name} ({t_range})  •  {t_tag}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = t_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 6: The Hero Case Study (WOW Moment)
    # ==========================================
    s6 = add_blank_slide()
    add_header(s6, "Live WOW Demonstration Case", "Hero Case Study: Intercepting ₹85,000 APP Scam")

    # Left: Victim & Context
    add_card(s6, 0.8, 1.8, 5.6, 5.0, border_color=COLOR_ROSE)
    tb = s6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Victim Profile & Coercion Context"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE

    details = [
        "Victim: Sunita Sharma (Age 72, Digital Skill 3/10)",
        "Transfer Request: ₹85,000 to Rahul Traders (PAY-6350)",
        "SMS Signal: 'SECURITY ALERT: Account frozen. Transfer immediately to safe vault account PAY-6350.'",
        "Vishing Call: 'Officer Sharma from Bank Security: Transfer funds to temporary safe account PAY-6350!'",
        "Authentication: User used real phone, correct PIN, valid OTP."
    ]
    for d in details:
        p = tf.add_paragraph()
        p.text = "• " + d
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(6)

    # Right: AI Verdict & Decision
    add_card(s6, 6.8, 1.8, 5.7, 5.0, border_color=COLOR_CYAN)
    tb = s6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "FinShield AI Interception Verdict"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    verdicts = [
        ("HUMAN MANIPULATION RISK", "94 / 100 (CRITICAL)"),
        ("Social NLP Coercion Score", "96.5 / 100"),
        ("Transaction Amount Ratio", "8.5x Normal Average"),
        ("Fraud Graph Reputation", "Connected to 7 Scam Networks"),
        ("ADAPTIVE FIREWALL DECISION", "⛔ TRANSACTION HELD | ₹85,000 SAVED")
    ]
    for lbl, val in verdicts:
        p = tf.add_paragraph()
        p.text = f"{lbl}: "
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_MUTED
        p.space_before = Pt(6)

        p_val = tf.add_paragraph()
        p_val.text = val
        p_val.font.size = Pt(14)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_ROSE if "HELD" in val or "94" in val else COLOR_WHITE

    # ==========================================
    # SLIDE 7: Accessibility & Senior Protection
    # ==========================================
    s7 = add_blank_slide()
    add_header(s7, "Inclusion & Accessibility", "Senior Safety & WCAG AAA Alignment")

    features = [
        ("WCAG 2.1 AAA Contrast", "High contrast white text (#FFFFFF) on slate-950 background (≥ 7:1 ratio) ensuring zero visual blur.", COLOR_CYAN),
        ("Color + Icon + Text Badges", "Multi-cue decision indicators complying with WCAG SC 1.4.1 (not relying on color cues alone).", COLOR_PURPLE),
        ("Speech Synthesis TTS", "Web Speech API auto-reads alert warnings at 0.85x speed for vision-impaired or senior users.", COLOR_AMBER),
        ("Family Trusted Contact", "Dispatches instant SMS alert to designated family guardian to double-check suspicious transfers.", COLOR_ROSE)
    ]

    for idx, (title, desc, col) in enumerate(features):
        row = idx // 2
        col_idx = idx % 2
        left = 0.8 + (col_idx * 5.95)
        top = 1.8 + (row * 2.6)
        
        add_card(s7, left, top, 5.7, 2.3, border_color=col)
        tb = s7.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(5.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 8: Quantum & Post-Quantum Security
    # ==========================================
    s8 = add_blank_slide()
    add_header(s8, "Future-Proof Cybersecurity", "Quantum ML & Post-Quantum Cryptography (PQC)")

    # Left: PQC
    add_card(s8, 0.8, 1.8, 5.6, 5.0, border_color=COLOR_CYAN)
    tb = s8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Post-Quantum Cryptography (PQC)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    pqc_items = [
        "NIST FIPS 203 ML-KEM-512 (Kyber) Key Encapsulation Mechanism.",
        "Protects banking session keys against future Shor's algorithm quantum decryption.",
        "Executed live handshake: 800-byte Public Key, 768-byte Ciphertext.",
        "256-bit AES-GCM Quantum-Resistant Session Shared Secret."
    ]
    for item in pqc_items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(8)

    # Right: Experimental Quantum ML
    add_card(s8, 6.8, 1.8, 5.7, 5.0, border_color=COLOR_PURPLE)
    tb = s8.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Quantum ML & QAOA Research"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE

    q_items = [
        "Qiskit 4-Qubit QSVC: ZZFeatureMap entangled quantum fidelity kernel classifier.",
        "Quantum Advantage Delta: +5.0% accuracy improvement over classical SVM on complex non-linear feature maps.",
        "QAOA QUBO Optimizer: Max-weight capacity solver prioritizing high-risk fraud cases under investigator constraints."
    ]
    for item in q_items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 9: Empirical Validation & Benchmark Metrics
    # ==========================================
    s9 = add_blank_slide()
    add_header(s9, "Performance Benchmarks", "Empirical Evaluation Metrics (10,000 Test Suite)")

    metrics = [
        ("94.2%", "Fraud Recall Sensitivity", "471 / 500 scams intercepted", COLOR_CYAN),
        ("3.8%", "False Intervention Rate", "Minimal friction on normal users", COLOR_PURPLE),
        ("42.5 ms", "API Decision Latency", "Sub-50ms real-time inline evaluation", COLOR_AMBER),
        ("₹18.45 L", "Protected Financial Value", "Prevented loss in synthetic test suite", COLOR_ROSE)
    ]

    for idx, (val, title, desc, col) in enumerate(metrics):
        row = idx // 2
        col_idx = idx % 2
        left = 0.8 + (col_idx * 5.95)
        top = 1.8 + (row * 2.6)

        add_card(s9, left, top, 5.7, 2.3, border_color=col)
        tb = s9.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(5.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(2)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 10: Conclusion & Hackathon Vision
    # ==========================================
    s10 = add_blank_slide()
    add_header(s10, "Summary & Business Vision", "FinShield Quantum: Next-Gen Banking Fraud Firewall")

    add_card(s10, 0.8, 1.8, 11.7, 5.0, border_color=COLOR_CYAN)
    tb = s10.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11.1), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Takeaways for Judges"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    takeaways = [
        "Solves the $1.5B+ APP Scam Gap: Detects victim intent manipulation rather than just technical transaction credentials.",
        "Multimodal AI Fusion: Integrates XGBoost, Z-score behavior, NLP coercion context, and NetworkX fraud graph intelligence.",
        "Enterprise Ops & Senior Safety: Dual experience for Fraud Ops and WCAG AAA senior accessibility with family alerts.",
        "Quantum-Ready Security: NIST ML-KEM-512 PQC session encryption paired with Qiskit QSVC and QAOA research.",
        "Empirically Proven: 94.2% Fraud Recall at 42.5ms sub-50ms API response latency."
    ]
    for t in takeaways:
        p = tf.add_paragraph()
        p.text = "✓ " + t
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(10)

    # Save output
    output_path = "FinShield_Quantum_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_deck()
